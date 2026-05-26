"""PresentationContext backed by python-pptx (offline file mode).

Works with .pptx files directly — no PowerPoint installation required
for most operations.  Image export is limited; use COM mode for that.
"""

from __future__ import annotations

import copy
import os
from typing import Any, List, Optional, Set, Tuple

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt, Inches

from .context import PresentationContext
from .proxies import BoundingBox, ColorInfo, ParagraphLocation, ShapeProxy, SlideProxy, TextRunLocation

# ── unit conversion constants ──────────────────────────────────
_POINTS_PER_EMU = 12700.0
_EMUS_PER_POINT = 12700


def _emu_to_points(emu: int) -> float:
    return round(emu / _POINTS_PER_EMU, 2)


def _points_to_emu(pts: float) -> int:
    return int(round(pts * _EMUS_PER_POINT))


def _points_to_emu_dim(pts: float) -> int:
    """Clamp to valid EMU range (0..~53 inches)."""
    val = _points_to_emu(pts)
    return max(0, min(val, 49151999))


class PptxPresentationContext(PresentationContext):
    """python-pptx based backend for offline .pptx file operations."""

    def __init__(self, filepath: Optional[str] = None):
        self._prs: Optional[Presentation] = None
        self._filepath: Optional[str] = None
        self._selected_slides: List[int] = []
        self._selected_shapes: List[ShapeProxy] = []
        if filepath:
            self._prs = Presentation(filepath)
            self._filepath = filepath

    # ── metadata ──────────────────────────────────────────────

    @property
    def name(self) -> str:
        if self._filepath:
            return os.path.splitext(os.path.basename(self._filepath))[0]
        return "Untitled"

    @property
    def filepath(self) -> Optional[str]:
        return self._filepath

    @property
    def slide_count(self) -> int:
        return len(self._prs.slides) if self._prs else 0

    @property
    def backend_type(self) -> str:
        return "pptx"

    # ── selection (simulated for file mode) ────────────────────

    def set_selected_slides(self, indices: List[int]) -> None:
        self._selected_slides = indices

    def set_selected_shapes(self, proxies: List[ShapeProxy]) -> None:
        self._selected_shapes = proxies

    def get_selected_shape_proxies(self) -> List[ShapeProxy]:
        return self._selected_shapes

    def get_selected_slide_indices(self) -> List[int]:
        return self._selected_slides

    # ── slides ────────────────────────────────────────────────

    def get_slide_proxies(self) -> List[SlideProxy]:
        result = []
        for i, slide in enumerate(self._prs.slides):
            result.append(SlideProxy(internal_ref=slide, index=i, name=slide.name or f"Slide {i+1}"))
        return result

    def get_all_shapes_on_slide(self, slide_index: int) -> List[ShapeProxy]:
        slide = self._prs.slides[slide_index]
        return self._shapes_from_slide(slide, slide_index)

    def get_all_shapes_all_slides(self) -> List[Tuple[int, ShapeProxy]]:
        result = []
        for i, slide in enumerate(self._prs.slides):
            for shape in self._shapes_from_slide(slide, i):
                result.append((i, shape))
        return result

    # ── shape geometry ────────────────────────────────────────

    def set_shape_position(self, shape_proxy: ShapeProxy, left: float, top: float) -> None:
        shape = shape_proxy.internal_ref
        shape.left = _points_to_emu_dim(left)
        shape.top = _points_to_emu_dim(top)
        shape_proxy.bounds.left = left
        shape_proxy.bounds.top = top

    def set_shape_size(self, shape_proxy: ShapeProxy, width: float, height: float) -> None:
        shape = shape_proxy.internal_ref
        shape.width = _points_to_emu_dim(width)
        shape.height = _points_to_emu_dim(height)
        shape_proxy.bounds.width = width
        shape_proxy.bounds.height = height

    def get_shape_geometry(self, shape_proxy: ShapeProxy) -> BoundingBox:
        shape = shape_proxy.internal_ref
        bb = BoundingBox(
            left=_emu_to_points(shape.left),
            top=_emu_to_points(shape.top),
            width=_emu_to_points(shape.width),
            height=_emu_to_points(shape.height),
        )
        shape_proxy.bounds = bb
        return bb

    # ── text ──────────────────────────────────────────────────

    def get_all_text_runs_all_slides(self) -> List[TextRunLocation]:
        result = []
        for si, slide in enumerate(self._prs.slides):
            for shape_proxy in self._shapes_from_slide(slide, si):
                if shape_proxy.has_text_frame:
                    shape = shape_proxy.internal_ref
                    for pi, para in enumerate(shape.text_frame.paragraphs):
                        for ri, _ in enumerate(para.runs):
                            result.append(TextRunLocation(shape_proxy, pi, ri))
        return result

    def get_text_runs_for_shapes(self, shape_proxies: List[ShapeProxy]) -> List[TextRunLocation]:
        result = []
        for sp in shape_proxies:
            if sp.has_text_frame:
                shape = sp.internal_ref
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    for ri, _ in enumerate(para.runs):
                        result.append(TextRunLocation(sp, pi, ri))
        return result

    def set_font_for_text_run(self, location: TextRunLocation, font_name: str,
                              font_size: Optional[float] = None) -> None:
        shape = location.shape_proxy.internal_ref
        run = shape.text_frame.paragraphs[location.paragraph_index].runs[location.run_index]
        run.font.name = font_name
        if font_size is not None:
            run.font.size = Pt(font_size)

    def set_paragraph_spacing(self, location: ParagraphLocation,
                              line_spacing: Optional[float] = None,
                              space_before: Optional[float] = None,
                              space_after: Optional[float] = None) -> None:
        shape = location.shape_proxy.internal_ref
        para = shape.text_frame.paragraphs[location.paragraph_index]
        if line_spacing is not None:
            para.line_spacing = Pt(line_spacing)
        if space_before is not None:
            para.space_before = Pt(space_before)
        if space_after is not None:
            para.space_after = Pt(space_after)

    # ── color ─────────────────────────────────────────────────

    def extract_slide_colors(self, slide_index: int) -> List[ColorInfo]:
        slide = self._prs.slides[slide_index]
        colors: List[ColorInfo] = []
        seen: Set[str] = set()
        for shape in slide.shapes:
            try:
                fill = shape.fill
                if fill and fill.fore_color and fill.fore_color.rgb:
                    hex_code = str(fill.fore_color.rgb)
                    if hex_code not in seen:
                        seen.add(hex_code)
                        colors.append(ColorInfo(hex_code=hex_code, source="shape fill"))
            except Exception:
                pass
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                hex_code = str(run.font.color.rgb)
                                if hex_code not in seen:
                                    seen.add(hex_code)
                                    colors.append(ColorInfo(hex_code=hex_code, source="text color"))
                        except Exception:
                            pass
        return colors

    # ── slide operations ──────────────────────────────────────

    def delete_slides(self, indices: List[int]) -> None:
        indices_sorted = sorted(indices, reverse=True)
        xml_slides = self._prs.part._element.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
        sld_id_list = list(xml_slides)
        for idx in indices_sorted:
            if 0 <= idx < len(sld_id_list):
                rId = sld_id_list[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rId:
                    self._prs.part.drop_rel(rId)
                sld_id_list[idx].getparent().remove(sld_id_list[idx])

    def move_slide(self, from_index: int, to_index: int) -> None:
        if from_index == to_index:
            return
        element = self._prs.part._element
        ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}sldId'
        sld_id_list = element.findall(f'.//{ns}')
        if 0 <= from_index < len(sld_id_list) and 0 <= to_index < len(sld_id_list):
            sld = sld_id_list.pop(from_index)
            sld_id_list.insert(to_index, sld)

    def duplicate_slide(self, index: int) -> None:
        slide = self._prs.slides[index]
        slide_layout = slide.slide_layout
        new_slide = self._prs.slides.add_slide(slide_layout)
        for shape in slide.shapes:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

    # ── image export ──────────────────────────────────────────

    def export_slide_as_image(self, slide_index: int, output_path: str,
                              fmt: str = "PNG", dpi: int = 150,
                              width: Optional[int] = None,
                              height: Optional[int] = None) -> None:
        """Export is limited in PPTX mode — raises NotImplementedError.
        Use COM context for slide-to-image export.
        """
        raise NotImplementedError(
            "Slide image export is only available when connected to a live PowerPoint session.\n"
            "Use [Connect to PowerPoint] for this feature."
        )

    # ── persistence ───────────────────────────────────────────

    def save(self) -> None:
        if self._filepath and self._prs:
            self._prs.save(self._filepath)

    def save_as(self, filepath: str) -> None:
        if self._prs:
            self._prs.save(filepath)
            self._filepath = filepath

    def close(self) -> None:
        self._prs = None
        self._filepath = None

    # ── internal helpers ──────────────────────────────────────

    def _shapes_from_slide(self, slide, slide_index: int) -> List[ShapeProxy]:
        result = []
        for shape in slide.shapes:
            try:
                has_text = shape.has_text_frame
            except Exception:
                has_text = False
            try:
                st = shape.shape_type
            except Exception:
                st = MSO_SHAPE_TYPE.AUTO_SHAPE
            result.append(ShapeProxy(
                internal_ref=shape,
                shape_type=int(st) if st else 0,
                name=shape.name,
                bounds=BoundingBox(
                    left=_emu_to_points(shape.left),
                    top=_emu_to_points(shape.top),
                    width=_emu_to_points(shape.width),
                    height=_emu_to_points(shape.height),
                ),
                rotation=shape.rotation,
                has_text_frame=has_text,
                slide_index=slide_index,
            ))
        return result
