"""WatermarkEngine — add/remove text or image watermarks."""

from typing import Optional

from src.backend.context import PresentationContext
from src.engine.base import BaseEngine, EngineResult


class WatermarkEngine(BaseEngine):
    WATERMARK_NAME_PREFIX = "_islide_watermark_"

    @staticmethod
    def add(
        context: PresentationContext,
        text: str = "Confidential",
        font_name: str = "Microsoft YaHei",
        font_size: float = 36.0,
        opacity: int = 30,
        rotation: float = -45.0,
        color_hex: str = "999999",
        image_path: Optional[str] = None,
    ) -> EngineResult:
        """Add watermark to all slides.

        Args:
            text: Watermark text (ignored if image_path is set)
            font_name: Font for text watermark
            font_size: Font size in points
            opacity: 0-100 transparency (0=invisible, 100=opaque)
            rotation: Rotation angle in degrees
            color_hex: Text color in hex
            image_path: Path to an image to use as watermark instead of text
        """
        try:
            r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
            alpha = opacity / 100.0

            added = 0
            for si in range(context.slide_count):
                if image_path:
                    added += WatermarkEngine._add_image_watermark(context, si, image_path, alpha, rotation)
                else:
                    added += WatermarkEngine._add_text_watermark(context, si, text, font_name, font_size, r, g, b, alpha, rotation)

            return EngineResult(
                success=True,
                message=f"Added watermark to {added} slides",
                data={"count": added}
            )
        except Exception as e:
            return EngineResult(success=False, message=str(e))

    @staticmethod
    def remove(context: PresentationContext) -> EngineResult:
        """Remove all watermarks added by this tool."""
        try:
            removed = 0
            for si, shape_proxy in context.get_all_shapes_all_slides():
                if shape_proxy.name.startswith(WatermarkEngine.WATERMARK_NAME_PREFIX):
                    try:
                        if context.backend_type == "com":
                            shape_proxy.internal_ref.Delete()
                        else:
                            sp = shape_proxy.internal_ref
                            sp.getparent().remove(sp)
                        removed += 1
                    except Exception:
                        pass

            return EngineResult(
                success=True,
                message=f"Removed {removed} watermarks",
                data={"count": removed}
            )
        except Exception as e:
            return EngineResult(success=False, message=str(e))

    @staticmethod
    def _add_text_watermark(context, slide_index, text, font_name, font_size, r, g, b, alpha, rotation):
        """Add a text watermark to a slide."""
        try:
            if context.backend_type == "com":
                slide = context._pres.Slides(slide_index + 1)
                shape = slide.Shapes.AddTextbox(1, 100, 100, 400, 60)
                shape.Name = f"{WatermarkEngine.WATERMARK_NAME_PREFIX}{slide_index}"
                shape.Rotation = rotation
                tf = shape.TextFrame
                tf.TextRange.Text = text
                tf.TextRange.Font.Name = font_name
                tf.TextRange.Font.Size = font_size
                tf.TextRange.Font.Color.RGB = (b << 16) | (g << 8) | r
                try:
                    shape.Fill.Transparency = 1.0 - alpha
                except Exception:
                    pass
            else:
                slide = context._prs.slides[slide_index]
                from pptx.util import Pt, Inches
                from pptx.enum.text import PP_ALIGN
                left = Inches(0.5)
                top_val = Inches(3)
                width = Inches(9)
                height = Inches(1.5)
                txBox = slide.shapes.add_textbox(left, top_val, width, height)
                txBox.name = f"{WatermarkEngine.WATERMARK_NAME_PREFIX}{slide_index}"
                txBox.rotation = rotation
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = text
                run.font.name = font_name
                run.font.size = Pt(font_size)
                from pptx.util import RGBColor
                run.font.color.rgb = RGBColor(r, g, b)
            return 1
        except Exception:
            return 0

    @staticmethod
    def _add_image_watermark(context, slide_index, image_path, alpha, rotation):
        """Add an image watermark to a slide."""
        try:
            if context.backend_type == "com":
                slide = context._pres.Slides(slide_index + 1)
                shape = slide.Shapes.AddPicture(image_path, False, True, 100, 100)
                shape.Name = f"{WatermarkEngine.WATERMARK_NAME_PREFIX}{slide_index}"
                shape.Rotation = rotation
            else:
                slide = context._prs.slides[slide_index]
                from pptx.util import Inches
                shape = slide.shapes.add_picture(image_path, Inches(1), Inches(1))
                shape.name = f"{WatermarkEngine.WATERMARK_NAME_PREFIX}{slide_index}"
                shape.rotation = rotation
            return 1
        except Exception:
            return 0
